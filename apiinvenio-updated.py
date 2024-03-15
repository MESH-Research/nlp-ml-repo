import time
import requests
import requests
import csv
import os
import fitz  # PyMuPDF
import docx
import pptx
import time
import requests

def get_records_and_download_files(api_url, endpoint, api_key, output_csv, download_path):
    try:
        authinfo = {'Authorization': f'Bearer {api_key}'}
        page = 1
        page_size = 100 
        has_more_records = True

        files_metadata = []

        with open(output_csv, 'w', newline='', encoding='utf-8') as file:
            writer = csv.writer(file)
            writer.writerow(['ID', 'Language', 'File Name'])

            while has_more_records:
                response = requests.get(f"{api_url}/{endpoint}?size={page_size}&page={page}", headers=authinfo)
                response.raise_for_status()
                data = response.json()

                for record in data['hits']['hits']:
                    record_id = record['id']
                    languages = ','.join([language['id'] for language in record['metadata'].get('languages', [])])
                    file_entries = record['files'].get('entries', {})

                    for file_name in file_entries:
                        writer.writerow([record_id, languages, file_name])
                        file_url = f"{api_url}/api/records/{record_id}/files/{file_name}/content"
                        local_file_path = os.path.join(download_path, file_name)
                        download_file(file_url, local_file_path, authinfo)

                        files_metadata.append({'id': record_id, 'metadata': record['metadata'], 'key': file_name})

                # Check if there are more records to fetch
                if len(data['hits']['hits']) < page_size:
                    has_more_records = False
                else:
                    page += 1

        print("Data written to CSV and files downloaded.")

        process_files(api_url, files_metadata, download_path, output_csv, authinfo)

    except requests.exceptions.HTTPError as e:
        print(f"HTTP error: {e}")
    except Exception as e:
        print(f"An error occurred: {e}")

def download_file(url, local_filename, authoinfo):
    max_retries = 5  # Maximum number of retries
    backoff_factor = 1  # Time to wait between retries, increases after each retry

    for attempt in range(max_retries):
        try:
            with requests.get(url, allow_redirects=True) as r:
                r.raise_for_status()
                x = requests.get(r.text, allow_redirects=True)
                filesize = x.headers['Content-length']
                # print(filesize)
                # print(x.status_code)
                with open(local_filename, 'wb') as f:
                    for chunk in x.iter_content(chunk_size=8192):
                        f.write(chunk)
                return local_filename
        except requests.exceptions.HTTPError as e:
            print(f"HTTP error: {e}")
            time.sleep(backoff_factor * (2 ** attempt))
        except Exception as e:
            print(f"Error downloading file: {e}")
            time.sleep(backoff_factor * (2 ** attempt))

    return None

def write_to_csv(filename, data):
    with open(filename, 'a', newline='', encoding='utf-8') as file:
        writer = csv.writer(file)
        writer.writerow(data)

def process_files(api_url, files_metadata, download_path, output_csv, authinfo):
    non_textual_files = []
    for file_data in files_metadata:
        file_id = file_data['id']
        file_language = ','.join([language['id'] for language in file_data['metadata'].get('languages', [])])
        file_name = file_data['key']

        file_url = f"{api_url}/api/records/{file_id}/files/{file_name}/content"
        local_file_path = os.path.join(download_path, file_name)
        download_file(file_url, local_file_path, authinfo)

        extracted_text = None
        if file_name.endswith('.pdf'):
            extracted_text = extract_text_from_pdf(local_file_path)
        elif file_name.endswith('.docx') or file_name.endswith('.doc'):
            extracted_text = extract_text_from_docx(local_file_path)
        elif file_name.endswith('.pptx'):
            extracted_text = extract_text_from_pptx(local_file_path)

        write_to_csv(output_csv, [file_id, file_language, file_name, extracted_text or ""])

        os.remove(local_file_path)

        if extracted_text is None:
            non_textual_files.append(file_name)

    with open("non_textual_files.txt", 'w') as f:
        f.write('\n'.join(non_textual_files))

#For PDF files
def extract_text_from_pdf(pdf_file_path):
    try:
        with fitz.open(pdf_file_path) as doc:
            text = ""
            for page_num in range(len(doc)):
                page = doc.load_page(page_num)
                text += page.get_text()
            return text
    except Exception as e:
        print(f"Error processing PDF:{e}")
        return None
    
#For docx files
def extract_text_from_docx(docx_file_path):
    try:
        doc = docx.Document(docx_file_path)
        text = ""
        for paragraph in doc.paragraphs:
            text += paragraph.text + "\n"
        return text
    except Exception as e:
        print(f"Error extracting text from Docx document: {e}")
        return None

#For pptx files
def extract_text_from_pptx(pptx_file_path):
    try:
        presentation = pptx.Presentation(pptx_file_path)
        text = ""
        for slide in presentation.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text += shape.text + "\n"
        return text
    except Exception as e:
        print(f"Error extracting text from PowerPoint file: {e}")
        return None

#Useful things
api_url = "https://invenio-dev.hcommons-staging.org"
api_key = "a3AEXr9qcGjzGGb9KjIZMNcZSCUCOzZFsIMDMzuwXjbKLHOGTzvuMkQiviP0"
api_endpoint = "api/records"
output_csv = "output.csv"
download_path = "download_files"

os.makedirs(download_path, exist_ok=True)
get_records_and_download_files(api_url, api_endpoint,api_key, output_csv, download_path)