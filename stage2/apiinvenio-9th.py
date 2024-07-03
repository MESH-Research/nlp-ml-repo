import os
import time
import logging
import requests
import csv
import sys
import textract
import fitz
import pytesseract
from PIL import Image
import docx
from docx import Document
import pptx
import subprocess
import speech_recognition as sr
from pydub import AudioSegment

#Make sure large text field in csv can be processed, this is useful when I examine the output.csv
csv.field_size_limit(sys.maxsize)

#Set up logging and make it display in my console
logging.basicConfig(filename='process.log', level=logging.INFO, format='%(asctime)s - %(levelname)s - %(message)s')
console_handler = logging.StreamHandler()
console_handler.setLevel(logging.INFO)
formatter = logging.Formatter('%(asctime)s - %(levelname)s - %(message)s')
console_handler.setFormatter(formatter)
logging.getLogger().addHandler(console_handler)

#I double checked here: defined function named "download_file", my goal is to download a file from a given 'url' and save it to 'local_filename')
def download_file(url, local_filename, authinfo):
    max_retries = 5  # Maximum number of retries
    backoff_factor = 1  # Time to wait between retries, increases after each retry

    for attempt in range(max_retries):
        try:
            with requests.get(url, headers= authinfo, allow_redirects=True) as r: #this makes an HTTP GET request to provided 'url'
                r.raise_for_status() #check for http error
                x = requests.get(r.text, allow_redirects=True) #this is my 2nd request, using the response text from the 1st request. Now im getting the actual file content.
                with open(local_filename, 'wb') as f: #content from 2nd request is written to the local file
                    for chunk in x.iter_content(chunk_size=8192): #read the data in chunks
                        f.write(chunk) #write to the local file
                return local_filename
        except requests.exceptions.HTTPError as e:
            logging.error(f"HTTP error on attempt: {attempt +1} for {url}: {e}")
        except Exception as e:
            logging.error(f"Error downloading file on attempt {attempt + 1} for {url}: {e}")
            
        time.sleep(backoff_factor * (2 ** attempt))

    return None

def extract_file(local_file_path, file_name):
    extraction_methods = [
        ('.pdf', extract_text_from_pdf),
        ('.docx', extract_text_from_docx),
        ('.doc', extract_text_from_doc),
        ('.pptx', extract_text_from_pptx),
        ('.csv', extract_text_from_csv),
        ('.txt', extract_text_from_txt),
        ('.mp3', extract_text_from_mp3)
    ]

    file_extension = os.path.splitext(file_name)[1].lower()

    for extension, func in extraction_methods:
        if file_extension == extension:
            try:
                return func(local_file_path)
            except Exception as e:
                logging.error(f"Error extracting text from {file_name} at {local_file_path}: {e}")
                return None
        
    logging.warning(f"Unsupported file extension for {file_name}. Skipping file.")
    return None

#For pdf, needs to decide if this file is text-based or image-based; if a, go through one lib; if b, go through another lib.
def extract_text_from_pdf(pdf_file_path):
    try:
        with fitz.open(pdf_file_path) as doc:
            text = ""
            for page_num in range(len(doc)):
                page = doc.load_page(page_num)
                page_text = page.get_text()
                if page_text.strip(): #If there is text on the page, it's likely text-based
                    text += page_text
                else: #This probably image-based
                    pix = page.get_pixmap()
                    img = Image.frombytes("RGB",[pix.width, pix.height], pix.samples)
                    text += pytesseract.image_to_string(img)
            return text
    except Exception as e:
        logging.error(f"Error processing PDF:{e}")
        return None
    
#For docx, it worked
def extract_text_from_docx(docx_file_path):
    try:
        doc = docx.Document(docx_file_path)
        text = ""
        for paragraph in doc.paragraphs:
            text += paragraph.text + "\n"
        return text
    except Exception as e:
        logging.error(f"Error extracting text from Docx document: {e}")
        return None
       
#For .doc, switch to docx, then extract
def convert_doc_to_docx(doc_file_path):
    try:
        output_dir = os.path.dirname(doc_file_path)
        base_file_name = os.path.splitext(base_file_name)[0]
        new_file_name = base_file_name + '.docx'
        newdocx_file_path = os.path.join(output_dir, new_file_name)

        result = subprocess.run(['libreoffice', '--headless', '--convert-to', '.docx','--outdir', output_dir, doc_file_path], check=True, capture_output=True)
        
        if result.returncode == 0:
            if os.path.exists(newdocx_file_path):
                return newdocx_file_path
            else:
                logging.error(f"Converted file not found at expected location: {newdocx_file_path}")
        else:
            logging.error(f"LibreOffice conversion failed: {result.stderr}")
            return None
        
    except subprocess.CalledProcessError as e:
        logging.error(f"Error converting DOC to DOCX: {e}")
        return None

def extract_text_from_doc(doc_file_path):
    try:
        if doc_file_path.endswith('.doc'):
            doc_file_path = convert_doc_to_docx(doc_file_path)

        doc = Document(doc_file_path)
        full_text = []
        for para in doc.paragraphs:
            full_text.append(para.text)
        return '\n'.join(full_text)
    except Exception as e:
        logging.error(f"Error extracting text from Docx document: {e}")
        return None

#For txt, check python builtins
def extract_text_from_txt(txt_file_path):
    try:
        with open(txt_file_path, 'r', encoding='utf-8') as file:
            text = file.read()
        return text
    except Exception as e:
        logging.error(f"Error extracting text from TXT file: {e}")

#For pptx, simple; I think it worked, with one exception
def extract_text_from_pptx(pptx_file_path):
    try:
        presentation = pptx.Presentation(pptx_file_path)
        text = ""
        for slide in presentation.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text += shape.text + "\n"
        return text
    except Exception as e:
        logging.error(f"Error extracting text from PowerPoint file: {e}")
        return None

#For csv, pythonbuiltins; it worked!
def extract_text_from_csv(csv_file_path):
    try:
        with open(csv_file_path, mode='r', encoding='utf-8') as file:
            reader = csv.reader(file)
            extracted_text = ''
            for row in reader:
                extracted_text += ''.join(row) + '\n'
            return extracted_text.strip()
    except Exception as e:
        logging.error(f"Error extracting text from CSV file: {e}")
        return None

def extract_text_from_mp3(local_file_path):
    recognizer = sr.Recognizer()
    try:
        sound = AudioSegment.from_mp3(local_file_path)
        wav_file_path = local_file_path.replace(".mp3", ".wav")
        sound.export(wav_file_path, format="wav")
    except Exception as e:
        logging.error(f"Error converting MP3 to WAV: {e}")
        return None

    try:
        with sr.AudioFile(wav_file_path) as source:
            audio_data = recognizer.record(source)
            text = recognizer.recognize_google(audio_data)
            return text
    except sr.UnknownValueError:
        logging.error("Speech could not be understood")
    except sr.RequestError as e:
        logging.error(f"Could not request results; {e}")
    finally:
        os.remove(wav_file_path)  # Delete the WAV file after processing

#This is the most complicated function here; it calls, download, extract, save, remove; or fail to extract, keep
def process_files(api_url, endpoint, api_key, output_csv, download_path, investigation_csv):
    try:
        authinfo = {'Authorization': f'Bearer {api_key}'}
        page = 1
        page_size = 100
        has_more_pages = True
        counter = 0

        # Headers for the CSV files
        output_headers = ['Record ID', 'Languages', 'File Name', 'Extracted Text']
        investigation_headers = ['Record ID', 'Languages', 'Files Name']

        while has_more_pages:
            response = requests.get(f"{api_url}/{endpoint}?size={page_size}&page={page}", headers=authinfo)
            response.raise_for_status()
            data = response.json()
            #This is where I check each page and all its records

            for record in data['hits']['hits']:
                counter += 1
                record_id = record['id']
                logging.info(f"record count {counter}, starting to process {record_id}")
                languages = ','.join([language['id']for language in record['metadata'].get('languages',[])])
                logging.info(f"record Id: {record_id}, language: {languages}")
 
                for file_name in record['files'].get('entries', {}):
                    local_file_path = os.path.join(download_path, file_name)
                    file_url = f"{api_url}/api/records/{record_id}/files/{file_name}/content"
        
                    if download_file(file_url, local_file_path, authinfo):
                        extracted_text = extract_file(local_file_path, file_name)

                        if not os.path.exists(output_csv):
                            with open(output_csv, 'w', newline='', encoding='utf-8') as file:
                                csv.writer(file).writerow(output_headers)
                        
                        if not os.path.exists(investigation_csv):
                            with open(investigation_csv, 'w', newline='', encoding='utf-8') as inv_file:
                                csv.writer(inv_file).writerow(investigation_headers)

                        with open(output_csv, 'a', newline='', encoding='utf-8') as file, \
                            open(investigation_csv, 'a', newline='', encoding='utf-8') as inv_file:
                            main_writer = csv.writer(file)
                            inv_writer = csv.writer(inv_file)

                            if extracted_text:
                                main_writer.writerow([record_id, languages, file_name, extracted_text])
                                logging.info(f"Processed and logged {record_id} under the name {file_name} in output.csv")
                                os.remove(local_file_path)
                            else:
                                inv_writer.writerow([record_id, languages, file_name])
                                logging.error(f"Failed to process {record_id} under the name {file_name}, logged in investigation.csv for further investigation")

            logging.info(f"Completed page {page}")

            if 'next' in data['links'].keys():                 
                page += 1   
            else:
                has_more_pages = False
            

        print("Data written to CSV and files processed.")

    except requests.exceptions.HTTPError as e:
        logging.error(f"HTTP error: {e}")
    except Exception as e:
        logging.error(f"An error occurred: {e}")

#Useful things
api_url = "https://invenio-dev.hcommons-staging.org"
api_key = "my api key"
api_endpoint = "api/records"
output_csv = "output9.csv"
investigation_csv = "investigation9.csv"
download_path = "download_files9"

os.makedirs(download_path, exist_ok=True)
process_files(api_url, api_endpoint,api_key, output_csv, download_path, investigation_csv)
